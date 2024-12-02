from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Sorting Techniques in Python"
subtitle.text = "An Overview of Common Sorting Algorithms"

# Bubble Sort slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Bubble Sort"
content.text = (
    "Bubble Sort is a simple sorting algorithm that repeatedly steps through the list, "
    "compares adjacent elements and swaps them if they are in the wrong order. "
    "The pass through the list is repeated until the list is sorted."
)

# Insertion Sort slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Insertion Sort"
content.text = (
    "Insertion Sort is a simple sorting algorithm that builds the final sorted array one item at a time. "
    "It is much less efficient on large lists than more advanced algorithms such as quicksort, heapsort, or merge sort."
)

# Merge Sort slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Merge Sort"
content.text = (
    "Merge Sort is an efficient, stable, comparison-based, divide and conquer sorting algorithm. "
    "Most implementations produce a stable sort, meaning that the implementation preserves the input order of equal elements in the sorted output."
)

# Quick Sort slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Quick Sort"
content.text = (
    "Quick Sort is an efficient sorting algorithm, serving as a systematic method for placing the elements of an array in order. "
    "Developed by Tony Hoare in 1959, it is still a commonly used algorithm for sorting."
)

from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Adding a shape with animation
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Animated Slide"

# Add a shape
left = Inches(2.0)
top = Inches(2.0)
width = Inches(4.0)
height = Inches(1.0)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = "This is an animated shape"

# Set the fill color
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

# Addfrom pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Adding a shape with animation
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Animated Slide"

# Add a shape
left = Inches(2.0)
top = Inches(2.0)
width = Inches(4.0)
height = Inches(1.0)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = "This is an animated shape"

# Set the fill color
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

# Add animation (this part needs to be done manually in PowerPoint) animation (this part needs to be done manually in PowerPoint)
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Adding a shape with animation
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Animated Slide"

# Add a shape
left = Inches(2.0)
top = Inches(2.0)
width = Inches(4.0)
height = Inches(1.0)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
shape.text = "This is an animated shape"

# Set the fill color
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

# Add animation (this part needs to be done manually in PowerPoint)
# Save the presentation
prs.save('sorting_techniques_in_python.pptx')
print("done")